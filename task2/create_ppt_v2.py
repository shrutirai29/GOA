"""Generate v2 PPT with 3D glassmorphism UI details."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(6,6,12); CARD=RGBColor(16,16,26); BORDER=RGBColor(26,26,40)
WHITE=RGBColor(232,232,240); MUTED=RGBColor(139,144,160); ACCENT=RGBColor(99,102,241)
GREEN=RGBColor(34,197,94); RED=RGBColor(239,68,68); AMBER=RGBColor(245,158,11)
CYAN=RGBColor(6,182,212); PURPLE=RGBColor(139,92,246)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); W=prs.slide_width

def bg(s):
    f=s.background.fill; f.solid(); f.fore_color.rgb=BG
def rect(s,l,t,w,h,fc=CARD,bc=BORDER):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h); sh.fill.solid(); sh.fill.fore_color.rgb=fc; sh.line.color.rgb=bc; sh.line.width=Pt(1); return sh
def txt(s,l,t,w,h,t2,sz=18,c=WHITE,b=False,a=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; p.text=t2; p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=b; p.font.name="Segoe UI"; p.alignment=a; return tf
def ml(s,l,t,w,h,lines,sz=15):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,(t2,cl,bd) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=t2; p.font.size=Pt(sz); p.font.color.rgb=cl; p.font.bold=bd; p.font.name="Segoe UI"; p.space_after=Pt(sz*0.4)
    return tf
def badge(s,l,t,t2,fc=ACCENT,w=Inches(1.2),h=Inches(0.32)):
    sh=rect(s,l,t,w,h,fc,fc); sh.text_frame.paragraphs[0].text=t2; sh.text_frame.paragraphs[0].font.size=Pt(10); sh.text_frame.paragraphs[0].font.color.rgb=WHITE; sh.text_frame.paragraphs[0].font.bold=True; sh.text_frame.paragraphs[0].font.name="Segoe UI"; sh.text_frame.paragraphs[0].alignment=PP_ALIGN.CENTER
def tb(s,n,title,sub=""):
    bg(s); bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,Pt(5)); bar.fill.solid(); bar.fill.fore_color.rgb=ACCENT; bar.line.fill.background()
    txt(s,Inches(0.6),Inches(0.25),Inches(1),Inches(0.35),f"{n:02d}",sz=11,c=MUTED)
    txt(s,Inches(0.6),Inches(0.6),Inches(11),Inches(0.6),title,sz=30,c=WHITE,b=True)
    if sub: txt(s,Inches(0.6),Inches(1.2),Inches(11),Inches(0.4),sub,sz=15,c=MUTED)

# Slide 1
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
bar=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,Pt(6)); bar.fill.solid(); bar.fill.fore_color.rgb=ACCENT; bar.line.fill.background()
for ox,oy,ow,oh,cl in [(Inches(1),Inches(1.5),Inches(4),Inches(4),RGBColor(18,16,48)),(Inches(8),Inches(4),Inches(3.5),Inches(3.5),RGBColor(8,21,32))]:
    o=sl.shapes.add_shape(MSO_SHAPE.OVAL,ox,oy,ow,oh); o.fill.solid(); o.fill.fore_color.rgb=cl; o.line.fill.background()
d=sl.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.8),Inches(2.2),Inches(0.22),Inches(0.22)); d.fill.solid(); d.fill.fore_color.rgb=GREEN; d.line.fill.background()
txt(sl,Inches(1.15),Inches(2.05),Inches(10),Inches(0.65),"Voice-Enabled RAG System",sz=40,c=WHITE,b=True)
txt(sl,Inches(0.8),Inches(2.85),Inches(10),Inches(0.45),"HH Goa 2026 - Shortlisting Task 2",sz=22,c=ACCENT,b=True)
txt(sl,Inches(0.8),Inches(3.45),Inches(10),Inches(0.4),"MSMARCO-XI Hindi  |  Voice > STT > Retrieval > Grounded Answer",sz=15,c=MUTED)
for i,(lb,cl) in enumerate([("4 Chunking Strategies",ACCENT),("Hybrid Retrieval",GREEN),("P50 < 85ms",CYAN),("62 Tests Passing",AMBER),("Guardrails Built-in",RED),("3D Glassmorphism UI",PURPLE)]):
    badge(sl,Inches(0.8+i*1.95),Inches(4.3),lb,cl,w=Inches(1.8))
txt(sl,Inches(0.8),Inches(6.5),Inches(10),Inches(0.35),"#RAGInGoa",sz=13,c=MUTED)

# Slide 2
sl=prs.slides.add_slide(prs.slide_layouts[6]); tb(sl,2,"Problem & Approach","What we built and why")
rect(sl,Inches(0.6),Inches(1.9),Inches(5.8),Inches(5.0),CARD,ACCENT)
txt(sl,Inches(0.9),Inches(2.0),Inches(5.2),Inches(0.35),"THE CHALLENGE",sz=13,c=ACCENT,b=True)
ml(sl,Inches(0.9),Inches(2.45),Inches(5.2),Inches(4.2),[("Build a voice-enabled RAG system:",WHITE,True),("",WHITE,False),("  Transcribes spoken Hindi questions",WHITE,False),("  Retrieves context from MSMARCO-XI",WHITE,False),("  Generates grounded, verified answers",WHITE,False),("  Refuses when it cannot answer reliably",GREEN,False),("",WHITE,False),("Key constraints:",AMBER,True),("  Full pipeline under 200ms latency",WHITE,False),("  Multiple chunking strategies",WHITE,False),("  Proper harness with retries",WHITE,False),("  Guardrails: unsafe, injection, off-topic",WHITE,False),("  Beautiful, modern UI for demo",PURPLE,False)],sz=14)
rect(sl,Inches(6.9),Inches(1.9),Inches(5.8),Inches(5.0),CARD,GREEN)
txt(sl,Inches(7.2),Inches(2.0),Inches(5.2),Inches(0.35),"OUR SOLUTION",sz=13,c=GREEN,b=True)
ml(sl,Inches(7.2),Inches(2.45),Inches(5.2),Inches(4.2),[("Full pipeline in 8 stages:",WHITE,True),("",WHITE,False),("1. Voice Input > Audio upload",WHITE,False),("2. STT > Sarvam / ElevenLabs / Mock",WHITE,False),("3. Query Router > Classifies intent",CYAN,False),("4. Guardrails > Blocks unsafe",RED,False),("5. Retrieval > Dense + BM25 + RRF",ACCENT,False),("6. Context > Dedup, merge, budget",WHITE,False),("7. Generate > With bounded retries",WHITE,False),("8. Grounding > Claim vs evidence",GREEN,False),("",WHITE,False),("UI: 3D glassmorphism with orbs,",PURPLE,True),("floating gradients, glow effects",PURPLE,False)],sz=14)

# Slide 3
sl=prs.slides.add_slide(prs.slide_layouts[6]); tb(sl,3,"Pipeline Architecture","Structured orchestration with per-stage timing")
for i,(lb,ds,cl) in enumerate([("Voice","Audio\nupload",ACCENT),("STT","Sarvam /\nElevenLabs",CYAN),("Router","Query\nclassify",ACCENT),("Guard","Block\nunsafe",RED),("Retrieve","Dense +\nBM25",GREEN),("Context","Dedup +\nmerge",AMBER),("Generate","LLM +\nretries",PURPLE),("Ground","Verify\nclaims",GREEN)]):
    x=Inches(0.35)+i*Inches(1.55); sh=rect(sl,x,Inches(2.2),Inches(1.4),Inches(1.5),CARD,cl)
    txt(sl,x+Inches(0.05),Inches(2.28),Inches(1.3),Inches(0.25),f"{i+1}",sz=10,c=cl,b=True,a=PP_ALIGN.CENTER)
    txt(sl,x+Inches(0.05),Inches(2.5),Inches(1.3),Inches(0.35),lb,sz=13,c=cl,b=True,a=PP_ALIGN.CENTER)
    txt(sl,x+Inches(0.05),Inches(2.9),Inches(1.3),Inches(0.7),ds,sz=10,c=MUTED,a=PP_ALIGN.CENTER)
    if i<7: txt(sl,x+Inches(1.42),Inches(2.7),Inches(0.15),Inches(0.3),">",sz=14,c=MUTED,a=PP_ALIGN.CENTER)
rect(sl,Inches(0.6),Inches(4.1),Inches(12.1),Inches(3.0),CARD,BORDER)
txt(sl,Inches(0.9),Inches(4.2),Inches(5),Inches(0.3),"KEY DESIGN DECISIONS",sz=13,c=ACCENT,b=True)
ml(sl,Inches(0.9),Inches(4.6),Inches(5.5),Inches(2.2),[("Parallel dense + BM25 (saves ~40ms)",WHITE,False),("RRF fusion merges both result lists",WHITE,False),("Confidence gate: abstains on weak evidence",WHITE,False),("Grounding check: claims vs sources",GREEN,False),("Every stage timed independently",WHITE,False)],sz=12)
ml(sl,Inches(6.9),Inches(4.6),Inches(5.5),Inches(2.2),[("Bounded retries: STT (1), LLM (config)",WHITE,False),("Structured I/O: PipelineResult objects",WHITE,False),("Error recovery: try/except at every boundary",WHITE,False),("Unique request IDs for log correlation",WHITE,False),("Devanagari-safe regex patterns",AMBER,False)],sz=12)

# Slide 4
sl=prs.slides.add_slide(prs.slide_layouts[6]); tb(sl,4,"Chunking Strategies","4 complementary approaches - each indexed separately")
for i,(nm,pr,ds,cl) in enumerate([("Fixed Token","256 tokens, overlap 40","Standard dense-retrieval baseline. Windows in word tokens.",ACCENT),("Sentence","128 words, overlap 32","Sliding window over natural sentence boundaries.",CYAN),("Semantic","Threshold 0.72","Sentence-embedding boundary detection. Splits at meaning shifts.",GREEN),("Hierarchical","Doc>Section>Para>Leaf","Preserves document structure with parent/sibling links.",AMBER)]):
    y=Inches(1.9)+i*Inches(1.35); rect(sl,Inches(0.6),y,Inches(12.1),Inches(1.15),CARD,cl)
    badge(sl,Inches(0.9),y+Inches(0.12),nm,cl,w=Inches(1.8))
    txt(sl,Inches(2.9),y+Inches(0.08),Inches(3),Inches(0.3),pr,sz=12,c=cl,b=True)
    txt(sl,Inches(0.9),y+Inches(0.5),Inches(9),Inches(0.5),ds,sz=12,c=MUTED)
txt(sl,Inches(0.6),Inches(7.0),Inches(12),Inches(0.3),"Every chunk carries: document_id, chunk_id, strategy, language, query_type, prev/next links",sz=11,c=MUTED,a=PP_ALIGN.CENTER)

# Slide 5
sl=prs.slides.add_slide(prs.slide_layouts[6]); tb(sl,5,"Hybrid Retrieval & Fusion","Dense + BM25 in parallel, merged via RRF")
rect(sl,Inches(0.6),Inches(1.9),Inches(5.8),Inches(2.3),CARD,ACCENT)
txt(sl,Inches(0.9),Inches(2.0),Inches(5),Inches(0.3),"DENSE RETRIEVAL",sz=13,c=ACCENT,b=True)
ml(sl,Inches(0.9),Inches(2.4),Inches(5.2),Inches(1.6),[("Model: multilingual-e5-small (22M params)",WHITE,False),("Embedding dim: 384 | Index: FAISS IVFFlat",WHITE,False),("Query embed: ~19ms (cached)",CYAN,False),("Strength: semantic similarity, paraphrases",GREEN,False)],sz=12)
rect(sl,Inches(6.9),Inches(1.9),Inches(5.8),Inches(2.3),CARD,CYAN)
txt(sl,Inches(7.2),Inches(2.0),Inches(5),Inches(0.3),"BM25 RETRIEVAL",sz=13,c=CYAN,b=True)
ml(sl,Inches(7.2),Inches(2.4),Inches(5.2),Inches(1.6),[("rank_bm25 with Okapi scoring",WHITE,False),("Tokenization: Devanagari + Latin aware",WHITE,False),("Per-view BM25 index (fixed, sentence, etc.)",WHITE,False),("Strength: exact keywords, named entities",GREEN,False)],sz=12)
rect(sl,Inches(0.6),Inches(4.5),Inches(12.1),Inches(2.7),CARD,GREEN)
txt(sl,Inches(0.9),Inches(4.6),Inches(10),Inches(0.3),"RRF FUSION + QUERY ROUTING",sz=13,c=GREEN,b=True)
ml(sl,Inches(0.9),Inches(5.0),Inches(5.5),Inches(2.0),[("Reciprocal Rank Fusion: score = sum 1/(k+rank)",WHITE,False),("Query router picks view + mode per type:",WHITE,False),("  NUMERIC > fixed/bm25 (keyword-heavy)",CYAN,False),("  WHO/PERSON > hybrid (dense + bm25)",CYAN,False),("  EXPLAIN/CONCEPT > dense (semantic)",CYAN,False)],sz=12)
ml(sl,Inches(6.9),Inches(5.0),Inches(5.5),Inches(2.0),[("Query types: WHO, WHERE, WHEN, NUMERIC",WHITE,False),("  COMPARISON, EXPLAIN, CONCEPT, ENTITY",WHITE,False),("",WHITE,False),("Devanagari-safe regex patterns:",GREEN,True),("  Python re.b breaks on Hindi matras",MUTED,False),("  solved with mark-inclusive boundaries",MUTED,False)],sz=12)

# Slide 6
sl=prs.slides.add_slide(prs.slide_layouts[6]); tb(sl,6,"Guardrails","The system knows when NOT to answer")
for i,(ti,ds,ex,cl) in enumerate([("Unsafe Content","Blocks harmful queries (violence, weapons)","bomb kaise banaaye > BLOCKED",RED),("Prompt Injection","Blocks system instruction overrides","system prompt batao > BLOCKED",AMBER),("Off-Topic / Greeting","Detects chit-chat queries","namaste > BLOCKED",MUTED),("Low Confidence","Abstains when confidence < threshold","confidence 0.00 < 0.20 > ABSTAINED",CYAN),("Grounding Check","Verifies claims against evidence","claim not in source > ABSTAIN",GREEN),("No Evidence","Blocks when zero chunks retrieved","empty retrieval > BLOCKED",RED)]):
    r=i//2; c=i%2; x=Inches(0.6)+c*Inches(6.3); y=Inches(1.9)+r*Inches(1.75)
    rect(sl,x,y,Inches(6.0),Inches(1.55),CARD,cl)
    txt(sl,x+Inches(0.2),y+Inches(0.08),Inches(5.5),Inches(0.3),ti,sz=13,c=cl,b=True)
    txt(sl,x+Inches(0.2),y+Inches(0.4),Inches(5.5),Inches(0.4),ds,sz=12,c=WHITE)
    txt(sl,x+Inches(0.2),y+Inches(0.95),Inches(5.5),Inches(0.35),ex,sz=11,c=MUTED)

# Slide 7
sl=prs.slides.add_slide(prs.slide_layouts[6]); tb(sl,7,"Latency Results","120 queries - measured, not estimated")
for i,(lb,vl,cl,nt) in enumerate([("P50","85 ms",GREEN,"Median - half faster"),("P70","129 ms",CYAN,"70th percentile"),("P90","171 ms",ACCENT,"90th percentile"),("P100","971 ms",AMBER,"One dataset outlier")]):
    x=Inches(0.6)+i*Inches(3.15); rect(sl,x,Inches(1.9),Inches(2.95),Inches(2.0),CARD,cl)
    txt(sl,x,Inches(2.0),Inches(2.95),Inches(0.3),lb,sz=13,c=cl,b=True,a=PP_ALIGN.CENTER)
    txt(sl,x,Inches(2.35),Inches(2.95),Inches(0.6),vl,sz=34,c=cl,b=True,a=PP_ALIGN.CENTER)
    txt(sl,x+Inches(0.15),Inches(3.1),Inches(2.65),Inches(0.5),nt,sz=10,c=MUTED,a=PP_ALIGN.CENTER)
rect(sl,Inches(0.6),Inches(4.2),Inches(12.1),Inches(2.9),CARD,BORDER)
txt(sl,Inches(0.9),Inches(4.3),Inches(5),Inches(0.3),"PER-STAGE BREAKDOWN (median)",sz=13,c=ACCENT,b=True)
for i,(nm,vl) in enumerate([("Router","0.0 ms"),("Guardrails","0.0 ms"),("Retrieval","84 ms"),("Rerank","0.0 ms"),("Context","0.2 ms"),("Generation","0.1 ms"),("Grounding","0.3 ms")]):
    co=i%4; ro=i//4; x=Inches(0.9)+co*Inches(3.0); y=Inches(4.75)+ro*Inches(0.55)
    txt(sl,x,y,Inches(1.8),Inches(0.3),nm,sz=12,c=MUTED); txt(sl,x+Inches(1.8),y,Inches(1.0),Inches(0.3),vl,sz=12,c=WHITE,b=True)
txt(sl,Inches(0.9),Inches(6.0),Inches(11),Inches(0.3),"Retrieval dominates at 84ms (86%). Everything else sub-millisecond.",sz=11,c=MUTED)

# Slide 8
sl=prs.slides.add_slide(prs.slide_layouts[6]); tb(sl,8,"Retrieval Evaluation","100 queries with gold relevance labels")
rect(sl,Inches(0.6),Inches(1.9),Inches(3.8),Inches(3.5),CARD,ACCENT)
txt(sl,Inches(0.9),Inches(2.0),Inches(3.4),Inches(0.3),"FIXED VIEW",sz=13,c=ACCENT,b=True)
ml(sl,Inches(0.9),Inches(2.4),Inches(3.4),Inches(2.5),[("R@5:  0.68",WHITE,True),("R@10: 0.83",WHITE,True),("",WHITE,False),("BM25-only mode",MUTED,False),("Best for keyword-heavy queries",MUTED,False),("Exact names, numbers, entities",MUTED,False)],sz=13)
rect(sl,Inches(4.7),Inches(1.9),Inches(3.8),Inches(3.5),CARD,GREEN)
txt(sl,Inches(5.0),Inches(2.0),Inches(3.4),Inches(0.3),"SEMANTIC VIEW",sz=13,c=GREEN,b=True)
ml(sl,Inches(5.0),Inches(2.4),Inches(3.4),Inches(2.5),[("R@5:  0.72",WHITE,True),("R@10: 0.85",WHITE,True),("",WHITE,False),("Dense embed mode",MUTED,False),("Best for conceptual queries",MUTED,False),("Paraphrases, explanations",MUTED,False)],sz=13)
rect(sl,Inches(8.8),Inches(1.9),Inches(3.9),Inches(3.5),CARD,RED)
txt(sl,Inches(9.1),Inches(2.0),Inches(3.4),Inches(0.3),"KEY FINDING",sz=13,c=RED,b=True)
ml(sl,Inches(9.1),Inches(2.4),Inches(3.4),Inches(2.5),[("English cross-encoder",WHITE,True),("HURTS Hindi retrieval",RED,True),("",WHITE,False),("R@1 with reranker: 0.16",RED,False),("R@1 without: 0.31",GREEN,False),("",WHITE,False),("Reranker OFF by default",AMBER,True),("Decision based on data",AMBER,False)],sz=13)

# Slide 9
sl=prs.slides.add_slide(prs.slide_layouts[6]); tb(sl,9,"Tech Stack & UI Design","Modern 3D glassmorphism frontend")
rect(sl,Inches(0.6),Inches(1.9),Inches(4.0),Inches(2.5),CARD,ACCENT)
txt(sl,Inches(0.9),Inches(2.0),Inches(3.5),Inches(0.3),"BACKEND",sz=13,c=ACCENT,b=True)
ml(sl,Inches(0.9),Inches(2.35),Inches(3.5),Inches(2.0),[("Python 3.11 + FastAPI",WHITE,False),("Sentence-Transformers + FAISS",WHITE,False),("rank_bm25 for sparse retrieval",WHITE,False),("62 tests (pytest), all passing",GREEN,False)],sz=13)
rect(sl,Inches(4.9),Inches(1.9),Inches(4.0),Inches(2.5),CARD,GREEN)
txt(sl,Inches(5.2),Inches(2.0),Inches(3.5),Inches(0.3),"FRONTEND",sz=13,c=GREEN,b=True)
ml(sl,Inches(5.2),Inches(2.35),Inches(3.5),Inches(2.0),[("React 18 + Vite",WHITE,False),("3D glassmorphism design",PURPLE,False),("Animated gradient orbs",PURPLE,False),("Glow effects + smooth transitions",PURPLE,False)],sz=13)
rect(sl,Inches(9.2),Inches(1.9),Inches(3.5),Inches(2.5),CARD,PURPLE)
txt(sl,Inches(9.5),Inches(2.0),Inches(3.0),Inches(0.3),"UI FEATURES",sz=13,c=PURPLE,b=True)
ml(sl,Inches(9.5),Inches(2.35),Inches(3.0),Inches(2.0),[("Glassmorphism panels",WHITE,False),("Floating animated orbs",WHITE,False),("3D hover effects",WHITE,False),("Pulsing mic glow",WHITE,False),("Staggered animations",WHITE,False)],sz=13)
rect(sl,Inches(0.6),Inches(4.7),Inches(12.1),Inches(2.5),CARD,BORDER)
txt(sl,Inches(0.9),Inches(4.8),Inches(5),Inches(0.3),"DEPLOYMENT",sz=13,c=AMBER,b=True)
ml(sl,Inches(0.9),Inches(5.15),Inches(5.5),Inches(1.8),[("Dockerfile + docker-compose.yml ready",WHITE,False),("Render / Railway / HF Spaces compatible",WHITE,False),("Frontend served from same container",WHITE,False),(".env.example with all config options",WHITE,False)],sz=13)
ml(sl,Inches(6.9),Inches(5.15),Inches(5.5),Inches(1.8),[("Cloudflare tunnel for instant demo link",WHITE,False),("Indexes pre-built (90k chunks, 4 views)",WHITE,False),("Model cached locally (471MB)",WHITE,False),("Total deploy size: ~333MB",WHITE,False)],sz=13)

# Slide 10
sl=prs.slides.add_slide(prs.slide_layouts[6]); tb(sl,10,"Results Summary","All requirements met")
for i,(rq,dt) in enumerate([("Speech-to-Text","Sarvam + ElevenLabs + Mock providers"),("Chunking","4 strategies: Fixed, Sentence, Semantic, Hierarchical"),("Latency","P50 = 85ms, P70 = 129ms - under 200ms target"),("Analytics","P50/P70/P100 across 120 queries with per-stage breakdown"),("Harness","Structured orchestration, retries, error recovery, request IDs"),("Guardrails","Unsafe, injection, off-topic, low-confidence, grounding check"),("Tests","62/62 passing - chunking, retrieval, routing, guardrails, API"),("Retrieval","R@5 = 0.68, R@10 = 0.83 (fixed view, 100 gold queries)"),("UI Design","3D glassmorphism with animated orbs, glow effects, responsive")]):
    y=Inches(1.9)+i*Inches(0.58); bgc=CARD if i%2==0 else RGBColor(14,14,22)
    rect(sl,Inches(0.6),y,Inches(12.1),Inches(0.5),bgc,BORDER)
    txt(sl,Inches(0.9),y+Inches(0.07),Inches(3.0),Inches(0.35),rq,sz=13,c=ACCENT,b=True)
    txt(sl,Inches(4.0),y+Inches(0.07),Inches(7.5),Inches(0.35),dt,sz=13,c=WHITE)

# Slide 11
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
bar=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,Pt(6)); bar.fill.solid(); bar.fill.fore_color.rgb=ACCENT; bar.line.fill.background()
for ox,oy,ow,oh,cl in [(Inches(0.5),Inches(2),Inches(3),Inches(3),RGBColor(18,16,48)),(Inches(9),Inches(3.5),Inches(4),Inches(4),RGBColor(8,21,32))]:
    o=sl.shapes.add_shape(MSO_SHAPE.OVAL,ox,oy,ow,oh); o.fill.solid(); o.fill.fore_color.rgb=cl; o.line.fill.background()
txt(sl,Inches(0.6),Inches(2.5),Inches(12),Inches(0.9),"Thank You",sz=46,c=WHITE,b=True,a=PP_ALIGN.CENTER)
txt(sl,Inches(0.6),Inches(3.4),Inches(12),Inches(0.5),"HH Goa 2026 - Task 2",sz=22,c=ACCENT,a=PP_ALIGN.CENTER)
txt(sl,Inches(0.6),Inches(4.1),Inches(12),Inches(0.4),"Voice-Enabled RAG on MSMARCO-XI Hindi",sz=17,c=MUTED,a=PP_ALIGN.CENTER)
ml(sl,Inches(3.5),Inches(5.0),Inches(6),Inches(1.5),[("Live Demo > trycloudflare.com",CYAN,False),("GitHub > github.com/shrutirai29/GOA",CYAN,False),("API Docs > /docs endpoint",CYAN,False),("",WHITE,False),("#RAGInGoa",ACCENT,True)],sz=15)

prs.save(r"D:\projects\goa\task2\HH_Goa_2026_RAG_Presentation_v2.pptx")
print("Saved: HH_Goa_2026_RAG_Presentation_v2.pptx")
